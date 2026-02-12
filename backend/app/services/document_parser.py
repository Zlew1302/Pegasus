"""Document parsing service — extracts text from various formats.

Adapted from the Orion project's DocumentParser.
Supports: PDF, DOCX, XLSX, PPTX, HTML, Markdown, JSON, Plain Text.
"""

import io
import re
import zipfile
from typing import Any


# Extension → file type mapping
EXTENSION_MAP: dict[str, str] = {
    "pdf": "pdf",
    "docx": "docx",
    "doc": "docx",
    "xlsx": "xlsx",
    "xls": "xlsx",
    "pptx": "pptx",
    "ppt": "pptx",
    "html": "html",
    "htm": "html",
    "md": "markdown",
    "markdown": "markdown",
    "json": "json",
    "txt": "text",
    "csv": "csv",
    "log": "text",
    "py": "text",
    "js": "text",
    "ts": "text",
    "png": "image",
    "jpg": "image",
    "jpeg": "image",
    "gif": "image",
    "webp": "image",
}


def detect_file_type(filename: str, content: bytes | None = None) -> str:
    """Detect file type from filename extension, optionally falling back to magic bytes."""
    if filename:
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
        if ext in EXTENSION_MAP:
            return EXTENSION_MAP[ext]

    if content:
        if content.startswith(b"%PDF"):
            return "pdf"
        elif content.startswith(b"PK"):
            return _detect_office_format(content)
        elif content.startswith(b"<!DOCTYPE") or content.startswith(b"<html"):
            return "html"
        elif content.startswith(b"{") or content.startswith(b"["):
            return "json"

    return "text"


def _detect_office_format(content: bytes) -> str:
    """Detect specific Office format from ZIP-based file."""
    try:
        with zipfile.ZipFile(io.BytesIO(content), "r") as zf:
            names = zf.namelist()
            if any(n.startswith("ppt/") for n in names):
                return "pptx"
            if any(n.startswith("xl/") for n in names):
                return "xlsx"
            if any(n.startswith("word/") for n in names):
                return "docx"
    except zipfile.BadZipFile:
        pass
    return "docx"


def parse_document(content: bytes, file_type: str) -> tuple[str, dict[str, Any]]:
    """Parse document and extract text content.

    Returns (extracted_text, metadata).
    """
    metadata: dict[str, Any] = {"original_type": file_type}

    if file_type == "pdf":
        text = _parse_pdf(content)
    elif file_type == "docx":
        text = _parse_docx(content)
    elif file_type == "xlsx":
        text = _parse_xlsx(content)
    elif file_type == "pptx":
        text = _parse_pptx(content)
    elif file_type == "html":
        text = _parse_html(content)
    elif file_type == "markdown":
        text = _parse_markdown(content)
    elif file_type == "csv":
        text = _parse_csv(content)
    elif file_type == "image":
        text = _parse_image(content)
    elif file_type == "json":
        text = _parse_json(content)
    else:
        text = content.decode("utf-8", errors="replace")

    text = _clean_text(text)
    metadata["character_count"] = len(text)
    metadata["word_count"] = len(text.split())

    return text, metadata


# ── Individual parsers ──────────────────────────────────────────


def _parse_pdf(content: bytes) -> str:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(content))
    parts: list[str] = []
    for page_num, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            parts.append(f"[Seite {page_num + 1}]\n{page_text}")
    return "\n\n".join(parts)


def _parse_docx(content: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(content))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name[-1] if para.style.name[-1].isdigit() else "1"
                parts.append(f"{'#' * int(level)} {para.text}")
            else:
                parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)


def _parse_xlsx(content: bytes) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(
        io.BytesIO(content), read_only=True, data_only=True
    )
    parts: list[str] = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        sheet_content = [f"## Sheet: {sheet_name}"]

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                sheet_content.append(" | ".join(row_values))

        if len(sheet_content) > 1:
            parts.append("\n".join(sheet_content))

    workbook.close()
    return "\n\n".join(parts)


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = [f"[Folie {slide_num}]"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_content.append(row_text)

        if len(slide_content) > 1:
            parts.append("\n".join(slide_content))

    return "\n\n".join(parts)


def _parse_html(content: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content.decode("utf-8", errors="replace"), "html.parser")
    for element in soup(["script", "style", "nav", "footer", "header"]):
        element.decompose()
    return soup.get_text(separator="\n")


def _parse_markdown(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _parse_csv(content: bytes) -> str:
    """Parse CSV content into a Markdown table (max 200 rows)."""
    import csv as csv_mod

    text = content.decode("utf-8", errors="replace")
    reader = csv_mod.reader(io.StringIO(text))
    rows: list[list[str]] = []
    for i, row in enumerate(reader):
        if i >= 201:  # 1 header + 200 data rows
            break
        rows.append(row)

    if not rows:
        return text

    # Build Markdown table
    header = rows[0]
    parts = [" | ".join(header)]
    parts.append(" | ".join("---" for _ in header))
    for row in rows[1:]:
        # Pad or trim to match header columns
        padded = row + [""] * (len(header) - len(row))
        parts.append(" | ".join(padded[: len(header)]))

    return "\n".join(parts)


def _parse_image(content: bytes) -> str:
    """Extract metadata from an image file."""
    size_kb = len(content) / 1024
    dimensions = ""

    try:
        from PIL import Image

        img = Image.open(io.BytesIO(content))
        fmt = img.format or "Unbekannt"
        w, h = img.size
        dimensions = f", {w}x{h}px"
    except (ImportError, Exception):
        # Detect format from magic bytes
        if content[:8] == b"\x89PNG\r\n\x1a\n":
            fmt = "PNG"
        elif content[:2] == b"\xff\xd8":
            fmt = "JPEG"
        elif content[:4] == b"GIF8":
            fmt = "GIF"
        elif content[:4] == b"RIFF" and content[8:12] == b"WEBP":
            fmt = "WEBP"
        else:
            fmt = "Bild"

    return f"[Bild: {fmt}{dimensions}, {size_kb:.0f} KB]"


def _parse_json(content: bytes) -> str:
    import json as json_mod

    data = json_mod.loads(content.decode("utf-8", errors="replace"))
    return _extract_text_from_json(data)


def _extract_text_from_json(obj: Any, depth: int = 0) -> str:
    if depth > 10:
        return ""
    if isinstance(obj, str):
        return obj
    elif isinstance(obj, dict):
        parts = []
        for key, value in obj.items():
            extracted = _extract_text_from_json(value, depth + 1)
            if extracted:
                parts.append(f"{key}: {extracted}")
        return "\n".join(parts)
    elif isinstance(obj, list):
        parts = [_extract_text_from_json(item, depth + 1) for item in obj]
        return "\n".join(p for p in parts if p)
    else:
        return str(obj) if obj is not None else ""


def _clean_text(text: str) -> str:
    """Normalize whitespace and line breaks."""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [line.strip() for line in text.split("\n")]
    return "\n".join(lines).strip()
